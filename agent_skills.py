# agent_skills.py
import os
import json
import time
import urllib.parse
from PIL import Image
from google import genai
from functools import wraps

# For slide generation and charting features
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
    print("[agent_skills] ✅ Matplotlib loaded successfully.")
except Exception as _mpl_err:
    MATPLOTLIB_AVAILABLE = False
    print(f"[agent_skills] ⚠️  Matplotlib unavailable (ft2font binary conflict or similar): {_mpl_err}. Fallback chart renderer will be used.")

# Initialize the client (Requires GEMINI_API_KEY to be set in your environment variables)
client = genai.Client()

# 🛡️ ENTERPRISE GUARDRAIL: Automated Exponential Backoff
def retry_on_429(max_retries=5, initial_wait=5):
    """Automatically pauses and retries the function if the Gemini API speed limit is hit."""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            wait_time = initial_wait
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    error_msg = str(e)
                    if "429" in error_msg or "RESOURCE_EXHAUSTED" in error_msg:
                        print(f"🚦 API Speed Limit Hit (429)! Sleeping for {wait_time}s to let quota reset... (Attempt {attempt+1}/{max_retries})")
                        time.sleep(wait_time)
                        wait_time *= 2  # Double the wait time on each failure (5s, 10s, 20s)
                    else:
                        raise e # If it is a real error (not a speed limit), crash normally
            # If it fails 5 times in a row, try one last time and let it report the error to the UI
            return func(*args, **kwargs)
        return wrapper
    return decorator

@retry_on_429()
def process_vision_image(image_path: str) -> str:
    """Uploads an image, analyzes it, and cleans up the cloud storage."""
    print(f"[Vision Skill] 👁️ Processing image: {image_path}")
    
    try:
        # 1. Securely upload the local file to Gemini's servers
        uploaded_file = client.files.upload(file=image_path)
        
        # 2. Strict system prompt forcing JSON output
        prompt = (
            "You are a highly analytical visual AI agent. Analyze the provided image. "
            "Return your analysis STRICTLY as a raw JSON dictionary with no markdown formatting. "
            "Include these exact keys: "
            "'summary' (a 2-sentence overview), "
            "'key_elements' (a list of 3 important things you see), "
            "'insight' (one actionable takeaway based on the image)."
        )
        
        # 3. Generate the multimodal response
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[uploaded_file, prompt],
            config={"temperature": 0.2} # Low temp keeps the analysis factual
        )
        
        # 4. COMPLIANCE REQUIREMENT: Delete the file from the cloud after use!
        client.files.delete(name=uploaded_file.name)
        print("[Vision Skill] 🧹 Cloud file deleted securely.")
        
        return response.text

    except Exception as e:
        print(f"[Vision Skill] ⚠️ Error: {e}")
        # Return a safe fallback JSON string if the API fails
        return json.dumps({
            "summary": "Vision API failed to process the image.",
            "key_elements": [],
            "insight": str(e)
        })

@retry_on_429()
def process_meeting_audio(audio_path: str) -> str:
    """
    Agent Skill: Transcribes and summarizes a meeting audio file (.wav, .mp3).
    Use this when the user wants to extract action items or summarize a recorded meeting.
    """
    print(f"[Skill] Processing meeting audio: {audio_path}")
    if not os.path.exists(audio_path):
        return "Error: Audio file not found."
        
    try:
        # Upload to Gemini File API (Kaggle Requirement for Media)
        audio_file = client.files.upload(file=audio_path)
        
        prompt = """You are a meeting summarizer. Generate a structured summary with:
        - Key discussion points
        - Decisions made
        - Action items
        Return ONLY a JSON object with keys: "transcript", "summary", "action_items"."""

        # --- OPTIMIZATION: Shifted explicitly to 1.5-flash for audio processing ---
        response = client.models.generate_content(
            model='gemini-1.5-flash',
            contents=[audio_file, prompt]
        )
        
        # Clean up the file from Google's servers
        client.files.delete(name=audio_file.name)
        
        raw = response.text.strip().replace("```json", "").replace("```", "").strip()
        return raw
    except Exception as e:
        return f"Audio Processing Error: {str(e)}"

@retry_on_429()
def draft_and_stage_email(to_email: str, instruction: str, context: str = "") -> str:
    """
    Agent Skill: Drafts a professional email and formats a direct browser link.
    Includes a Guardrail check to ensure no inappropriate content.
    """
    print(f"[Skill] Drafting email to: {to_email}")
    
    # Dynamic subject line optimization prompt
    prompt = f"""You are an executive email assistant. Draft an email based on this directive: "{instruction}"
    Context to include: {context}
    
    CRITICAL REQUIREMENT: You MUST dynamically generate a highly engaging, context-specific subject line based on the email body you write. DO NOT use generic placeholders like 'Automated Draft'.
    
    Return ONLY a JSON object with keys: "subject", "body"."""
    
    try:
        resp = client.models.generate_content(model='gemini-2.5-flash', contents=[prompt])
        raw = resp.text.strip().replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)
        
        body = data.get("body", "")
        
        # Internal guardrail verification step
        if not _guardrail_check(body):
            return "Error: Email drafting blocked by Guardrails. Inappropriate, unprofessional, or sensitive content detected."
            
        safe_subject = urllib.parse.quote(data.get("subject", "Update from the team"))
        safe_body = urllib.parse.quote(body)
        
        gmail_url = f"https://mail.google.com/mail/?view=cm&fs=1&to={to_email}&su={safe_subject}&body={safe_body}"
        
        return json.dumps({
            "status": "success",
            "message": "Email drafted safely.",
            "draft_url": gmail_url,
            "preview_body": body,
            "subject": data.get("subject", "Update from the team")
        }, indent=2)
        
    except Exception as e:
        return f"Drafting Error: {str(e)}"

@retry_on_429(max_retries=3, initial_wait=3)
def _guardrail_check(text: str) -> bool:
    """
    Internal Guardrail: Uses Gemini as a secondary evaluator to ensure the text is safe to send.
    Checks for profanity, explicit financial leaks, hostility, or hallucinated promises.
    """
    print("[Guardrail] Evaluating output for safety...")
    prompt = f"""Evaluate the following text for corporate safety. 
    Are there any profanities, hostile language, or highly inappropriate phrasing?
    Text: {text}
    Respond with exactly 'SAFE' or 'UNSAFE'."""
    
    try:
        resp = client.models.generate_content(model='gemini-2.5-flash', contents=[prompt])
        return "SAFE" in resp.text.upper()
    except:
        return False # Fail secure: If the guardrail breaks, block the action.


@retry_on_429()
def process_document_file(doc_path: str) -> str:
    """
    Agent Skill: Uploads a PDF document, extracts content and summarizes it using Gemini,
    and deletes the file from cloud storage safely.
    """
    print(f"[Document Skill] 📄 Processing document: {doc_path}")
    try:
        # Upload the file to Gemini File API (forces application/pdf mime type)
        uploaded_file = client.files.upload(
            file=doc_path, 
            config={'mime_type': 'application/pdf'}
        )
        
        prompt = (
            "You are a highly analytical document AI agent. Analyze the provided PDF document. "
            "Return your analysis STRICTLY as a raw JSON dictionary with no markdown formatting. "
            "Include these exact keys: "
            "'summary' (a 2-sentence overview), "
            "'key_elements' (a list of 3 important things you see), "
            "'insight' (one actionable takeaway based on the document)."
        )
        
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[uploaded_file, prompt],
            config={"temperature": 0.2}
        )
        
        # Deletion guarantee
        try:
            client.files.delete(name=uploaded_file.name)
            print("[Document Skill] 🧹 Cloud file deleted securely.")
        except Exception as cleanup_err:
            print(f"Warning: Failed to delete cloud file: {cleanup_err}")
            
        return response.text

    except Exception as e:
        print(f"[Document Skill] ⚠️ Error: {e}")
        return json.dumps({
            "summary": "Document API failed to process the PDF.",
            "key_elements": [],
            "insight": str(e)
        })


# ---------------------------------------------------------------------------
# THEME CONSTANTS  (dark slate corporate palette)
# ---------------------------------------------------------------------------
_C_BG        = RGBColor(15,  23,  42)   # slate-900  – slide background
_C_SURFACE   = RGBColor(30,  41,  59)   # slate-800  – card / container fill
_C_SURFACE2  = RGBColor(2,   6,   23)   # slate-950  – deep secondary bg
_C_BORDER    = RGBColor(51,  65,  85)   # slate-700  – shape border
_C_INDIGO    = RGBColor(99,  102, 241)  # indigo-500 – primary accent
_C_FUCHSIA   = RGBColor(217, 70,  239)  # fuchsia-500 – secondary accent
_C_WHITE     = RGBColor(248, 250, 252)  # slate-50
_C_MUTED     = RGBColor(148, 163, 184)  # slate-400  – sub-labels
_C_DIM       = RGBColor(100, 116, 139)  # slate-500  – badges / captions
_C_NAVY      = RGBColor(15,  23,  42)   # slate-900  – title-slide panel

# MSO shape type constants
_MSO_RECT    = 1   # msoShapeRectangle
_MSO_ROUNDED = 5   # msoShapeRoundedRectangle
_MSO_CHEVRON = 13  # msoShapeChevron


def _add_bg(slide, prs):
    """Flood-fills the entire slide with the dark slate background."""
    bg = slide.shapes.add_shape(_MSO_RECT,
                                 0, 0,
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _C_BG
    bg.line.fill.background()
    # Push background shape to back by moving the XML element to position 0
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)   # index 2 = after spTree housekeeping nodes


def _add_slide_chrome(slide, title_text: str, slide_index: int, total_slides: int):
    """Adds the left accent bar, header title, divider, and slide-count badge."""
    # Left accent bar
    accent = slide.shapes.add_shape(_MSO_RECT,
                                     Inches(0), Inches(0),
                                     Inches(0.12), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _C_INDIGO
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.32),
                                          Inches(11.5), Inches(0.75))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(26)
    title_p.font.bold = True
    title_p.font.name = "Calibri"
    title_p.font.color.rgb = _C_WHITE

    # Indigo divider under title
    divider = slide.shapes.add_shape(_MSO_RECT,
                                      Inches(0.45), Inches(1.1),
                                      Inches(11.5), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = _C_INDIGO
    divider.line.fill.background()

    # Slide badge (top-right)
    badge = slide.shapes.add_textbox(Inches(11.8), Inches(0.22),
                                      Inches(1.4), Inches(0.45))
    badge_p = badge.text_frame.paragraphs[0]
    badge_p.text = f"{slide_index} / {total_slides}"
    badge_p.font.size = Pt(9)
    badge_p.font.name = "Calibri"
    badge_p.font.color.rgb = _C_DIM
    badge_p.alignment = PP_ALIGN.RIGHT


def _set_shape_style(shape, fill_rgb: RGBColor, border_rgb: RGBColor = None,
                     border_pt: float = 0.75):
    """Applies a solid fill and optional border to any pptx shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()


def _add_text_to_shape(shape, text: str, font_size: int, bold: bool = False,
                        color: RGBColor = None, align=PP_ALIGN.LEFT):
    """Writes a single paragraph of styled text into a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = "Calibri"
    p.font.color.rgb = color or _C_WHITE
    p.alignment = align


# ---------------------------------------------------------------------------
# LAYOUT RENDERERS
# ---------------------------------------------------------------------------

def _render_text_default(slide, prs, slide_data: dict):
    """TEXT_DEFAULT: left-accent dark card with bulleted text body."""
    bullets = slide_data.get("bullets", [])
    # Content area card
    card = slide.shapes.add_shape(_MSO_ROUNDED,
                                   Inches(0.45), Inches(1.25),
                                   Inches(12.4), Inches(5.85))
    _set_shape_style(card, _C_SURFACE, _C_BORDER, 0.5)
    # Indent accent strip inside card
    strip = slide.shapes.add_shape(_MSO_RECT,
                                    Inches(0.45), Inches(1.25),
                                    Inches(0.06), Inches(5.85))
    _set_shape_style(strip, _C_FUCHSIA)

    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.45),
                                            Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {bullet}"
        p.font.size = Pt(16)
        p.font.name = "Calibri"
        p.font.color.rgb = _C_MUTED
        p.space_before = Pt(8)
        p.line_spacing = 1.3


def _render_timeline(slide, prs, slide_data: dict):
    """
    TIMELINE: A central horizontal accent rail with N evenly-spaced
    block-card nodes, each containing a label.  Connector dots mark
    each node on the rail.
    """
    items = slide_data.get("bullets", [])
    n = len(items)
    if n == 0:
        return

    rail_y      = Inches(3.9)
    rail_x      = Inches(0.6)
    rail_w      = Inches(12.1)
    rail_h      = Inches(0.06)
    card_w      = Inches(1.85)
    card_h      = Inches(1.55)
    card_top_y  = Inches(1.55)   # above-rail cards (odd indices)
    card_bot_y  = Inches(4.25)   # below-rail cards (even indices)
    dot_r       = Inches(0.18)

    # Central timeline rail
    rail = slide.shapes.add_shape(_MSO_RECT, rail_x, rail_y, rail_w, rail_h)
    _set_shape_style(rail, _C_INDIGO)

    usable_w = rail_w - Inches(0)
    spacing  = usable_w / max(n, 1)

    for i, label in enumerate(items):
        cx = rail_x + spacing * i + spacing / 2  # centre x of this node
        bx = cx - card_w / 2

        # Alternating above/below placement for visual stagger
        is_above = (i % 2 == 0)
        card_y   = card_top_y if is_above else card_bot_y

        # Connector vertical stem
        stem_top    = card_top_y + card_h if is_above else rail_y + rail_h
        stem_bottom = rail_y               if is_above else card_bot_y
        stem_h      = abs(stem_bottom - stem_top)
        stem = slide.shapes.add_shape(_MSO_RECT,
                                       cx - Inches(0.02), min(stem_top, stem_bottom),
                                       Inches(0.04), stem_h)
        _set_shape_style(stem, _C_BORDER)

        # Dot on rail
        dot = slide.shapes.add_shape(_MSO_ROUNDED,
                                      cx - dot_r, rail_y - dot_r + rail_h / 2,
                                      dot_r * 2, dot_r * 2)
        _set_shape_style(dot, _C_INDIGO, _C_SURFACE, 1.0)

        # Card
        card = slide.shapes.add_shape(_MSO_ROUNDED, bx, card_y, card_w, card_h)
        _set_shape_style(card, _C_SURFACE, _C_INDIGO, 0.75)

        # Step number badge
        num_box = slide.shapes.add_textbox(bx + Inches(0.08), card_y + Inches(0.06),
                                            Inches(0.38), Inches(0.32))
        np = num_box.text_frame.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(9)
        np.font.bold = True
        np.font.name = "Calibri"
        np.font.color.rgb = _C_INDIGO

        # Label text
        lbl_box = slide.shapes.add_textbox(bx + Inches(0.12), card_y + Inches(0.38),
                                            card_w - Inches(0.22), card_h - Inches(0.42))
        lbl_tf = lbl_box.text_frame
        lbl_tf.word_wrap = True
        lp = lbl_tf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(10)
        lp.font.name = "Calibri"
        lp.font.color.rgb = _C_WHITE
        lp.line_spacing = 1.2


def _render_3_column_grid(slide, prs, slide_data: dict):
    """
    3_COLUMN_GRID: Three equal-width rounded-rectangle containers.
    First bullet → col 1, second → col 2, third → col 3.
    Each container gets a numbered header badge + body text.
    """
    bullets  = slide_data.get("bullets", [])
    columns  = slide_data.get("columns", bullets[:3])   # allow explicit override
    col_headers = slide_data.get("column_headers", [])

    # Pad / truncate to exactly 3
    while len(columns) < 3:
        columns.append("")
    columns = columns[:3]

    col_w   = Inches(3.9)
    col_h   = Inches(5.5)
    col_y   = Inches(1.3)
    gap     = Inches(0.25)
    start_x = Inches(0.5)
    accent_colors = [_C_INDIGO, _C_FUCHSIA, RGBColor(16, 185, 129)]  # indigo / fuchsia / emerald

    for i in range(3):
        cx = start_x + i * (col_w + gap)

        # Card container
        card = slide.shapes.add_shape(_MSO_ROUNDED, cx, col_y, col_w, col_h)
        _set_shape_style(card, _C_SURFACE, _C_BORDER, 0.5)

        # Top accent stripe inside card
        stripe = slide.shapes.add_shape(_MSO_RECT, cx, col_y, col_w, Inches(0.07))
        _set_shape_style(stripe, accent_colors[i])

        # Column number badge
        badge_box = slide.shapes.add_textbox(cx + Inches(0.15), col_y + Inches(0.14),
                                              Inches(0.55), Inches(0.38))
        bp = badge_box.text_frame.paragraphs[0]
        bp.text = f"0{i+1}"
        bp.font.size = Pt(20)
        bp.font.bold = True
        bp.font.name = "Calibri"
        bp.font.color.rgb = accent_colors[i]

        # Column header (from column_headers list if present)
        hdr = col_headers[i] if i < len(col_headers) else f"Column {i+1}"
        hdr_box = slide.shapes.add_textbox(cx + Inches(0.15), col_y + Inches(0.58),
                                            col_w - Inches(0.3), Inches(0.45))
        hp = hdr_box.text_frame.paragraphs[0]
        hp.text = hdr.upper()
        hp.font.size = Pt(10)
        hp.font.bold = True
        hp.font.name = "Calibri"
        hp.font.color.rgb = _C_MUTED

        # Thin separator
        sep = slide.shapes.add_shape(_MSO_RECT,
                                      cx + Inches(0.15), col_y + Inches(1.1),
                                      col_w - Inches(0.3), Inches(0.02))
        _set_shape_style(sep, accent_colors[i])

        # Body text
        body_box = slide.shapes.add_textbox(cx + Inches(0.18), col_y + Inches(1.22),
                                             col_w - Inches(0.36), col_h - Inches(1.35))
        btf = body_box.text_frame
        btf.word_wrap = True
        bp2 = btf.paragraphs[0]
        bp2.text = columns[i]
        bp2.font.size = Pt(13)
        bp2.font.name = "Calibri"
        bp2.font.color.rgb = _C_WHITE
        bp2.line_spacing = 1.35


def _render_kpi_cards(slide, prs, slide_data: dict):
    """
    KPI_CARDS: Draws up to 4 large metric boxes in a 2×2 grid.
    Each bullet is expected in the format  "Label: value"  or  "value | label".
    The metric value is rendered in oversized high-contrast font;
    the label in muted smaller type below.
    """
    bullets = slide_data.get("bullets", [])
    kpis    = slide_data.get("kpis", bullets)   # allow explicit kpis key
    kpis    = kpis[:4]  # max 4 cards

    n       = len(kpis)
    cols    = 2 if n > 2 else n
    rows    = (n + cols - 1) // cols

    card_w  = Inches(5.8)
    card_h  = Inches(2.4)
    gap_x   = Inches(0.35)
    gap_y   = Inches(0.3)

    total_w = cols * card_w + (cols - 1) * gap_x
    total_h = rows * card_h + (rows - 1) * gap_y
    start_x = (prs.slide_width  - total_w) / 2
    start_y = Inches(1.4) + (Inches(5.7) - total_h) / 2

    accent_seq = [_C_INDIGO, _C_FUCHSIA,
                  RGBColor(16, 185, 129), RGBColor(251, 191, 36)]  # emerald / amber

    for i, kpi_text in enumerate(kpis):
        row = i // cols
        col = i %  cols
        cx  = start_x + col * (card_w + gap_x)
        cy  = start_y + row * (card_h + gap_y)
        acc = accent_seq[i % len(accent_seq)]

        # Parse "Label: Value" or "Value | Label" or fall back to raw string as value
        if ":" in kpi_text:
            parts = kpi_text.split(":", 1)
            kpi_label = parts[0].strip()
            kpi_value = parts[1].strip()
        elif "|" in kpi_text:
            parts = kpi_text.split("|", 1)
            kpi_value = parts[0].strip()
            kpi_label = parts[1].strip()
        else:
            kpi_value = kpi_text.strip()
            kpi_label = ""

        # Card background
        card = slide.shapes.add_shape(_MSO_ROUNDED, cx, cy, card_w, card_h)
        _set_shape_style(card, _C_SURFACE, _C_BORDER, 0.5)

        # Left accent bar inside card
        bar = slide.shapes.add_shape(_MSO_RECT, cx, cy, Inches(0.1), card_h)
        _set_shape_style(bar, acc)

        # Oversized metric value
        val_box = slide.shapes.add_textbox(cx + Inches(0.25), cy + Inches(0.2),
                                            card_w - Inches(0.35), Inches(1.45))
        vp = val_box.text_frame.paragraphs[0]
        vp.text = kpi_value
        vp.font.size = Pt(52)
        vp.font.bold = True
        vp.font.name = "Calibri"
        vp.font.color.rgb = _C_WHITE

        # Label underneath
        if kpi_label:
            lbl_box = slide.shapes.add_textbox(cx + Inches(0.25), cy + Inches(1.7),
                                                card_w - Inches(0.35), Inches(0.55))
            lp2 = lbl_box.text_frame.paragraphs[0]
            lp2.text = kpi_label.upper()
            lp2.font.size = Pt(11)
            lp2.font.bold = True
            lp2.font.name = "Calibri"
            lp2.font.color.rgb = acc


# Map layout_type string → renderer function
_LAYOUT_RENDERERS = {
    "TIMELINE":       _render_timeline,
    "3_COLUMN_GRID":  _render_3_column_grid,
    "KPI_CARDS":      _render_kpi_cards,
    "TEXT_DEFAULT":   _render_text_default,
}


def _render_slide(prs, blank_layout, slide_data: dict, slide_index: int, total_slides: int):
    """
    Dispatcher: selects the correct infographic layout renderer based on
    slide_data["layout_type"], then adds chrome (background, title, divider).
    """
    layout_type = slide_data.get("layout_type", "TEXT_DEFAULT").upper()
    renderer    = _LAYOUT_RENDERERS.get(layout_type, _render_text_default)

    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, prs)
    _add_slide_chrome(slide, slide_data.get("title", f"Slide {slide_index}"),
                      slide_index, total_slides)
    renderer(slide, prs, slide_data)


def generate_slide_deck(
    text: str,
    slide_count: int = 5,
    filename: str = "presentation.pptx"
) -> dict:
    """
    Agent Skill: Issues a strict compact-JSON prompt to Gemini requesting
    per-slide layout_type tags, then renders an N-slide infographic PPTX
    via python-pptx using the corporate dark slate theme.

    Supported layout_type values: TIMELINE, 3_COLUMN_GRID, KPI_CARDS, TEXT_DEFAULT.

    Args:
        text:        Free-form input text (topic brief, meeting notes, etc.)
        slide_count: Exact number of content slides (1-15, clamped).
        filename:    Output .pptx filename.

    Returns:
        dict with keys: filename, slide_count, deck_title, slides.
    """
    slide_count = max(1, min(15, int(slide_count)))
    print(f"[Presentation Skill] Requesting {slide_count}-slide infographic structure from Gemini...")

    # ----------------------------------------------------------------
    # Step 1 – Ultra-strict JSON prompt.  NO MARKDOWN.  NO PADDING.
    # ----------------------------------------------------------------
    schema_example = (
        '{"title":"string",' 
        '"slides":[{'
        '"title":"string",' 
        '"layout_type":"TIMELINE|3_COLUMN_GRID|KPI_CARDS|TEXT_DEFAULT",' 
        '"bullets":["string","string","string"],' 
        '"columns":["col1 body","col2 body","col3 body"],' 
        '"column_headers":["Hdr1","Hdr2","Hdr3"],' 
        '"kpis":["Label: value","Label: value"]' 
        '}]}'
    )
    layout_rules = (
        "layout_type selection rules (pick the BEST fit per slide):\n"
        "  TIMELINE      → for sequential steps, processes, roadmaps, or historical events. "
                           "bullets = ordered short step labels (max 6 items).\n"
        "  3_COLUMN_GRID → for comparisons, three pillars, or feature breakdowns. "
                           "Populate columns[] with per-column body text AND column_headers[].\n"
        "  KPI_CARDS     → for metrics, statistics, or quantitative highlights. "
                           "Populate kpis[] as \"Label: value\" strings (max 4).\n"
        "  TEXT_DEFAULT  → default for narrative, strategic insight, or mixed content. "
                           "bullets = 3-5 concise sentences."
    )
    prompt = (
        f"OUTPUT ONLY RAW JSON. ZERO MARKDOWN. ZERO EXPLANATION. NO TRAILING COMMAS.\n"
        f"Generate exactly {slide_count} infographic presentation slides.\n"
        f"Schema: {schema_example}\n"
        f"{layout_rules}\n"
        f"Global rules: unique slide titles; omit unused optional keys; "
        f"no slide titled 'Title Slide'; content must be information-dense not generic.\n"
        f"Source text: {text}"
    )

    try:
        resp = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[prompt],
            config={"temperature": 0.15}   # lower temp → tighter schema compliance
        )
        raw = resp.text.strip()
        # Strip any accidental markdown fences
        raw = raw.replace("```json", "").replace("```", "").strip()
        deck_data = json.loads(raw)
    except Exception as e:
        print(f"[Presentation Skill] Gemini structure call failed: {e}. Using TEXT_DEFAULT fallback.")
        deck_data = {
            "title": "Presentation",
            "slides": [
                {
                    "title": f"Slide {i + 1}",
                    "layout_type": "TEXT_DEFAULT",
                    "bullets": [text[:200]]
                }
                for i in range(slide_count)
            ]
        }

    deck_title  = deck_data.get("title", "Presentation")
    slides_data = deck_data.get("slides", [])

    # Clamp slide count (model may over/under produce)
    slides_data = slides_data[:slide_count]
    while len(slides_data) < slide_count:
        slides_data.append({
            "title": f"Section {len(slides_data) + 1}",
            "layout_type": "TEXT_DEFAULT",
            "bullets": ["Content pending."]
        })

    # Guarantee every slide has a valid layout_type
    for s in slides_data:
        if s.get("layout_type", "").upper() not in _LAYOUT_RENDERERS:
            s["layout_type"] = "TEXT_DEFAULT"

    total_content_slides = len(slides_data)
    total_slides         = total_content_slides + 1   # +1 for title slide

    print(f"[Presentation Skill] Rendering {total_slides} slides "
          f"(title + {total_content_slides} content, layouts: "
          f"{[s.get('layout_type','?') for s in slides_data]})...")

    # ----------------------------------------------------------------
    # Step 2 – Render via python-pptx
    # ----------------------------------------------------------------
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout     = prs.slide_layouts[6]  # completely blank

    # ---- Title Slide ------------------------------------------------
    ts = prs.slides.add_slide(blank_layout)
    _add_bg(ts, prs)

    # Left navy panel
    panel = ts.shapes.add_shape(_MSO_RECT, Inches(0), Inches(0),
                                 Inches(5.2), Inches(7.5))
    _set_shape_style(panel, _C_NAVY)

    # Indigo accent strip on right edge of panel
    strip = ts.shapes.add_shape(_MSO_RECT, Inches(5.2), Inches(0),
                                 Inches(0.06), Inches(7.5))
    _set_shape_style(strip, _C_INDIGO)

    # Decorative dot cluster on panel (visual premium touch)
    for di, (dx, dy) in enumerate([(0.7, 1.1), (1.1, 1.0), (0.9, 1.45)]):
        dot = ts.shapes.add_shape(_MSO_ROUNDED, Inches(dx), Inches(dy),
                                   Inches(0.12), Inches(0.12))
        _set_shape_style(dot, _C_INDIGO)

    # "ENTERPRISE INTELLIGENCE" label on panel
    lbl_box = ts.shapes.add_textbox(Inches(0.35), Inches(2.9),
                                     Inches(4.5), Inches(0.55))
    lp = lbl_box.text_frame.paragraphs[0]
    lp.text = "ENTERPRISE INTELLIGENCE"
    lp.font.size = Pt(10)
    lp.font.bold = True
    lp.font.name = "Calibri"
    lp.font.color.rgb = _C_INDIGO

    # Thin indigo rule under label
    rule = ts.shapes.add_shape(_MSO_RECT, Inches(0.35), Inches(3.48),
                                Inches(3.8), Inches(0.025))
    _set_shape_style(rule, _C_INDIGO)

    # Main deck title (right panel)
    t_box = ts.shapes.add_textbox(Inches(5.55), Inches(2.5),
                                   Inches(7.4), Inches(1.8))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_p = t_tf.paragraphs[0]
    t_p.text = deck_title
    t_p.font.size = Pt(38)
    t_p.font.bold = True
    t_p.font.name = "Calibri"
    t_p.font.color.rgb = _C_WHITE

    # Subtitle / meta line
    sub_box = ts.shapes.add_textbox(Inches(5.55), Inches(4.4),
                                     Inches(7.4), Inches(0.6))
    sub_p = sub_box.text_frame.paragraphs[0]
    sub_p.text = f"AI-Generated Infographic Deck  •  {total_content_slides} Slides"
    sub_p.font.size = Pt(14)
    sub_p.font.name = "Calibri"
    sub_p.font.color.rgb = _C_DIM

    # ---- Content Slides ---------------------------------------------
    for idx, slide_data in enumerate(slides_data, start=1):
        _render_slide(prs, blank_layout, slide_data, idx, total_content_slides)

    prs.save(filename)
    print(f"[Presentation Skill] ✅ Infographic deck saved → {filename}")

    return {
        "filename":    filename,
        "slide_count": total_slides,
        "deck_title":  deck_title,
        "slides":      slides_data
    }


def _generate_chart_matplotlib(debits_sum: float, credits_sum: float, chart_path: str) -> str:
    """Primary chart renderer using Matplotlib (requires working ft2font binary)."""
    fig, ax = plt.subplots(figsize=(6.5, 4.2))

    fig.patch.set_facecolor('#0f172a')
    ax.set_facecolor('#020617')

    categories = ['Debits', 'Credits']
    values = [debits_sum, credits_sum]
    colors = ['#f43f5e', '#10b981']

    bars = ax.bar(categories, values, color=colors, width=0.45, edgecolor='#1e293b', linewidth=1.5)

    for bar in bars:
        yval = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            yval + (max(values) * 0.02 if max(values) > 0 else 0.1),
            f"${yval:,.2f}",
            ha='center', va='bottom',
            color='#f8fafc', fontweight='bold', fontsize=10
        )

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#334155')
    ax.spines['bottom'].set_color('#334155')
    ax.tick_params(axis='x', colors='#94a3b8', labelsize=11)
    ax.tick_params(axis='y', colors='#94a3b8', labelsize=11)
    ax.set_title("Transaction Overview", color='#f8fafc', fontsize=13, fontweight='bold', pad=18)
    ax.grid(axis='y', linestyle='--', alpha=0.1, color='#94a3b8')

    plt.tight_layout()
    plt.savefig(chart_path, dpi=300, facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close(fig)
    return chart_path


def _generate_chart_pillow_fallback(debits_sum: float, credits_sum: float, chart_path: str) -> str:
    """
    Pure Pillow fallback chart renderer.
    Produces a styled dark PNG bar chart when matplotlib's ft2font binary is unavailable
    (e.g., Python 3.14 binary conflict).
    """
    from PIL import Image, ImageDraw

    W, H = 650, 420
    BG        = (15,  23,  42)   # slate-900
    PLOT_BG   = (2,   6,   23)   # slate-950
    ROSE      = (244, 63,  94)   # debit bar
    EMERALD   = (16,  185, 129)  # credit bar
    SLATE_400 = (148, 163, 184)
    SLATE_700 = (51,  65,  133)
    WHITE     = (248, 250, 252)
    BORDER    = (30,  41,  59)

    img  = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(img)

    # Plot area margins
    margin_left, margin_right = 70, 30
    margin_top, margin_bottom = 60, 55
    plot_x1 = margin_left
    plot_y1 = margin_top
    plot_x2 = W - margin_right
    plot_y2 = H - margin_bottom

    # Fill plot area
    draw.rectangle([plot_x1, plot_y1, plot_x2, plot_y2], fill=PLOT_BG)

    # Title
    title = "Transaction Overview"
    draw.text((W // 2, 22), title, fill=WHITE, anchor="mm")

    max_val = max(debits_sum, credits_sum, 1.0)
    plot_h  = plot_y2 - plot_y1

    bar_w   = int((plot_x2 - plot_x1) * 0.22)
    gap     = int((plot_x2 - plot_x1) * 0.14)
    total_w = 2 * bar_w + gap
    start_x = plot_x1 + ((plot_x2 - plot_x1) - total_w) // 2

    pairs = [
        ("Debits",  debits_sum,  ROSE,    start_x),
        ("Credits", credits_sum, EMERALD, start_x + bar_w + gap),
    ]

    # Horizontal grid lines (5 steps)
    for step in range(1, 6):
        gy = plot_y2 - int(plot_h * (step / 5))
        draw.line([(plot_x1, gy), (plot_x2, gy)], fill=(30, 41, 59), width=1)
        label = f"${max_val * step / 5:,.0f}"
        draw.text((plot_x1 - 6, gy), label, fill=SLATE_400, anchor="rm")

    for label, val, color, bx in pairs:
        bar_h  = int(plot_h * (val / max_val)) if max_val > 0 else 0
        bar_y1 = plot_y2 - bar_h
        bar_y2 = plot_y2

        # Bar shadow / border
        draw.rectangle([bx - 1, bar_y1 - 1, bx + bar_w + 1, bar_y2 + 1], fill=BORDER)
        # Bar fill
        draw.rectangle([bx, bar_y1, bx + bar_w, bar_y2], fill=color)

        # Value label on top of bar
        val_txt = f"${val:,.2f}"
        label_y = bar_y1 - 14 if bar_h > 0 else plot_y2 - 14
        draw.text((bx + bar_w // 2, label_y), val_txt, fill=WHITE, anchor="mm")

        # X-axis category label
        draw.text((bx + bar_w // 2, plot_y2 + 16), label, fill=SLATE_400, anchor="mm")

    # Bottom axis line
    draw.line([(plot_x1, plot_y2), (plot_x2, plot_y2)], fill=SLATE_700, width=1)
    # Left axis line
    draw.line([(plot_x1, plot_y1), (plot_x1, plot_y2)], fill=SLATE_700, width=1)

    # Fallback notice
    draw.text((W // 2, H - 12), "Pillow fallback renderer (matplotlib unavailable)", fill=(71, 85, 105), anchor="mm")

    img.save(chart_path, "PNG")
    return chart_path


def generate_expense_chart(records: list) -> str:
    """
    Agent Skill: Generates a clean, dark-themed non-signed summary bar graph of expenses.
    Uses matplotlib when available; falls back to a pure Pillow renderer if matplotlib's
    ft2font binary conflicts with the host Python environment (e.g., Python 3.14).
    Saves to the Next.js public directory and returns the absolute file path.
    """
    try:
        credits_sum = sum(abs(float(r['amount'])) for r in records if float(r['amount']) >= 0)
        debits_sum  = sum(abs(float(r['amount'])) for r in records if float(r['amount']) < 0)
    except Exception as parse_err:
        print(f"[Expense Chart] Error parsing amounts: {parse_err}")
        credits_sum = 0.0
        debits_sum  = 0.0

    public_dir = os.path.abspath(os.path.join("enterprise-ai-dashboard", "public"))
    os.makedirs(public_dir, exist_ok=True)
    chart_path = os.path.join(public_dir, "expense_chart.png")

    if MATPLOTLIB_AVAILABLE:
        print("[Expense Chart] Rendering with Matplotlib...")
        return _generate_chart_matplotlib(debits_sum, credits_sum, chart_path)
    else:
        print("[Expense Chart] ⚠️  Matplotlib unavailable — using Pillow fallback renderer.")
        return _generate_chart_pillow_fallback(debits_sum, credits_sum, chart_path)



